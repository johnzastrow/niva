#!/usr/bin/env python3
"""Build the niva pitch deck (PPTX) — CEO sales narrative for decision-makers and
non-programmer data scientists (Marimo users). White background, niva brand palette,
native editable text + embedded diagram graphics (rendered from presentation/graphics/).

Run:  uv run --no-project --with python-pptx python build_deck.py
Out:  presentation/niva_pitch.pptx
"""

import os
import struct

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
GFX = os.path.join(HERE, "graphics", "png")
LOGOS = os.path.join(os.path.dirname(HERE), "logos")
OUT = os.path.join(HERE, "niva_pitch.pptx")

# ---- brand palette -------------------------------------------------------
INK = RGBColor(0x28, 0x32, 0x3B)      # slate — headings / primary text
DEEP = RGBColor(0x46, 0x57, 0x53)     # deep sage — primary accent
SAGE = RGBColor(0x6B, 0x7B, 0x6F)     # mid sage — secondary text
MIST = RGBColor(0x92, 0x98, 0x8E)     # muted
TAN = RGBColor(0xC1, 0xB4, 0x9B)      # warm accent / highlight
PAIN = RGBColor(0xB8, 0x5C, 0x3C)     # terracotta — "without niva"
LIGHT = RGBColor(0xEE, 0xF1, 0xEC)    # light sage fill
PAINBG = RGBColor(0xFB, 0xF3, 0xEF)   # light terracotta fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = "Arial"

EMU_PER_IN = 914400
SW, SH = 13.333, 7.5  # 16:9


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def _bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _box(slide, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _fill(sp, color):
    sp.fill.solid()
    sp.fill.fore_color.rgb = color


def _text(slide, x, y, w, h, runs, size=18, color=INK, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # `runs` is a string or list of paragraphs; each paragraph a string or list of (text, opts)
    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        chunks = para if isinstance(para, list) else [(para, {})]
        for txt, opt in chunks:
            r = p.add_run()
            r.text = txt
            r.font.name = opt.get("font", font)
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.italic = opt.get("italic", italic)
            r.font.color.rgb = opt.get("color", color)
    return tb


def _footer(slide, n):
    # thin accent rule
    rule = _box(slide, 0.0, 7.18, SW, 0.06)
    _fill(rule, TAN)
    _text(slide, 0.55, 7.04, 6, 0.4,
          [[("niva", {"bold": True, "color": DEEP}), ("  ·  Easy wins every time", {"color": DEEP})]],
          size=12)
    _text(slide, SW - 1.55, 7.04, 1.0, 0.4, str(n), size=12, color=DEEP, bold=True, align=PP_ALIGN.RIGHT)


def content_slide(prs, title, eyebrow=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_white(s)
    # accent tick + title
    tick = _box(s, 0.55, 0.62, 0.12, 0.62)
    _fill(tick, DEEP)
    ty = 0.55
    if eyebrow:
        _text(s, 0.8, 0.5, 11.8, 0.35, eyebrow.upper(),
              size=12.5, color=TAN, bold=True)
        ty = 0.82
    _text(s, 0.8, ty, 11.9, 0.9, title, size=30, color=INK, bold=True)
    return s


def fit_image(slide, path, x, y, w, h, align="center"):
    iw, ih = png_size(path)
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        draw_w = w
        draw_h = w / ar
    else:
        draw_h = h
        draw_w = h * ar
    dx = x + (w - draw_w) / 2
    dy = y + (h - draw_h) / 2
    slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(draw_w), Inches(draw_h))


def image_slide(prs, title, img, eyebrow=None, caption=None):
    s = content_slide(prs, title, eyebrow)
    top = 1.7
    bottom = 6.95 if not caption else 6.6
    fit_image(s, os.path.join(GFX, img), 0.55, top, SW - 1.1, bottom - top)
    if caption:
        _text(s, 0.8, 6.62, 11.8, 0.4, caption, size=14, color=SAGE, italic=False,
              align=PP_ALIGN.CENTER)
    return s


def bullets_slide(prs, title, items, eyebrow=None, lead=None):
    s = content_slide(prs, title, eyebrow)
    y = 1.75
    if lead:
        _text(s, 0.85, y, 11.6, 0.9, lead, size=18, color=SAGE)
        y += 0.95
    for head, sub in items:
        dot = _box(s, 0.9, y + 0.07, 0.16, 0.16)
        _fill(dot, DEEP)
        _text(s, 1.25, y - 0.06, 11.2, 0.9,
              [[(head + "  ", {"bold": True, "color": INK, "size": 19}),
                (sub, {"color": SAGE, "size": 17})]], size=19)
        y += 0.86
    return s


# =========================================================================
prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_PER_IN))
prs.slide_height = Emu(int(SH * EMU_PER_IN))

# ---- 1. TITLE -----------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg_white(s)
logo = os.path.join(LOGOS, "logo_text.png")
lw, lh = png_size(logo)
disp_h = 2.7
disp_w = disp_h * lw / lh
s.shapes.add_picture(logo, Inches((SW - disp_w) / 2), Inches(0.7), Inches(disp_w), Inches(disp_h))
_text(s, 1, 3.65, SW - 2, 0.7, "Readable geoprocessing for everyone",
      size=30, color=INK, bold=True, align=PP_ALIGN.CENTER)
_text(s, 1, 4.45, SW - 2, 0.5,
      "Spatial analysis that a non-programmer can write, reproduce, and share.",
      size=18, color=SAGE, align=PP_ALIGN.CENTER)
band = _box(s, (SW - 3.4) / 2, 5.35, 3.4, 0.6)
_fill(band, INK)
_text(s, (SW - 3.4) / 2, 5.42, 3.4, 0.5, "Easy wins every time", size=18, color=TAN,
      bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
_text(s, 1, 6.5, SW - 2, 0.4, "Decision-maker & analyst briefing", size=13,
      color=MIST, align=PP_ALIGN.CENTER)

# ---- 2. THE OPPORTUNITY (statement) ------------------------------------
s = content_slide(prs, "Spatial questions are everywhere.", eyebrow="The opportunity")
_text(s, 0.85, 2.1, 11.6, 2.2,
      [[("The people who need the answers — analysts, scientists, planners — ",
         {"size": 26, "color": INK}),
        ("mostly aren't programmers.", {"size": 26, "color": DEEP, "bold": True})]],
      line_spacing=1.15)
_text(s, 0.85, 4.4, 11.6, 1.6,
      [[("Yet to ", {"size": 20, "color": SAGE}),
        ("automate", {"size": 20, "color": INK, "bold": True}),
        (" the work — to make it repeatable and shareable — today's tools demand "
         "that they write code. So the work waits, or doesn't happen.", {"size": 20, "color": SAGE})]],
      line_spacing=1.2)
_footer(s, 2)

# ---- 3. THE PROBLEM -----------------------------------------------------
s = bullets_slide(
    prs, "Automating QGIS means writing PyQGIS", eyebrow="The problem",
    lead="To script even simple geoprocessing in QGIS today, you must:",
    items=[
        ("Initialize a QgsApplication", "boilerplate before any work begins"),
        ("Memorize algorithm IDs", "like native:buffer, native:dissolve, native:clip"),
        ("Build ALL_CAPS parameter dicts", "and juggle TEMPORARY_OUTPUT between steps"),
        ("Thread each tool's output into the next", "by hand, step after step"),
    ])
_text(s, 1.25, 6.25, 11.2, 0.6,
      [[("That's a ", {"size": 19, "color": INK}),
        ("programming task", {"size": 19, "color": PAIN, "bold": True}),
        (" — out of reach for the GUI-first analysts who need it most, "
         "and tedious even for those who can.", {"size": 19, "color": INK})]])
_footer(s, 3)

# ---- 4. BEFORE / AFTER (graphic) ---------------------------------------
image_slide(prs, "Automation shouldn't require a programmer",
            "before_after.png", eyebrow="The shift")
_footer(prs.slides[-1], 4)

# ---- 5. MEET NIVA -------------------------------------------------------
s = content_slide(prs, "Meet niva", eyebrow="The product")
_text(s, 0.85, 1.7, 11.6, 1.0,
      [[("A concise, readable ", {"size": 24, "color": INK}),
        ("text-pipeline grammar", {"size": 24, "color": DEEP, "bold": True}),
        (" for QGIS geoprocessing — for people who don't want to write PyQGIS.",
         {"size": 24, "color": INK})]], line_spacing=1.15)
code = _box(s, 0.85, 3.15, 11.6, 1.25)
_fill(code, INK)
_text(s, 1.2, 3.15, 11.0, 1.25,
      [[("load", {"color": TAN, "bold": True}), (" roads.gpkg ", {"color": WHITE}),
        ("| ", {"color": MIST}), ("buffer", {"color": TAN, "bold": True}), (" 100 dissolve ", {"color": WHITE}),
        ("| ", {"color": MIST}), ("clip", {"color": TAN, "bold": True}), (" city.gpkg ", {"color": WHITE}),
        ("| ", {"color": MIST}), ("save", {"color": TAN, "bold": True}), (" roads_local.gpkg", {"color": WHITE})]],
      size=19, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
_text(s, 0.85, 4.75, 11.6, 1.2,
      [[("A whole pipeline on one line — that a non-programmer can ", {"size": 20, "color": SAGE}),
        ("write and read", {"size": 20, "color": INK, "bold": True}),
        (" — running on QGIS's own Processing algorithms underneath.", {"size": 20, "color": SAGE})]],
      line_spacing=1.2)
_text(s, 0.85, 6.1, 11.6, 0.5, "Easy wins every time.", size=18, color=TAN, italic=True, bold=True)
_footer(s, 5)

# ---- 6. WHO IT'S FOR ----------------------------------------------------
s = content_slide(prs, "Who niva is for", eyebrow="Audience")
cards = [
    ("Decision-makers", "Faster answers without hiring scarce GIS programmers. Work that's auditable, reproducible, and shareable across the org.", DEEP),
    ("Data scientists\n(non-programmers)", "Already use Marimo for repeatable, shareable analysis. Now add geoprocessing — in the same readable, reactive workflow.", TAN),
    ("Downstream consumers", "Explore the results in QGIS, in Marimo, or as static files — whichever fits, no GIS expertise required.", DEEP),
]
cw, gap = 3.85, 0.18
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (h, body, accent) in enumerate(cards):
    x = x0 + i * (cw + gap)
    card = _box(s, x, 2.0, cw, 4.4)
    _fill(card, LIGHT)
    top = _box(s, x, 2.0, cw, 0.22)
    _fill(top, accent)
    _text(s, x + 0.25, 2.45, cw - 0.5, 1.3, h, size=21, color=INK, bold=True, line_spacing=1.0)
    _text(s, x + 0.25, 3.75, cw - 0.5, 2.4, body, size=16, color=SAGE, line_spacing=1.15)
_footer(s, 6)

# ---- 7. FITS YOUR MARIMO WORKFLOW --------------------------------------
s = bullets_slide(
    prs, "It fits the workflow you already have", eyebrow="Marimo-native",
    lead="Marimo is how your team makes analysis repeatable and shareable. niva runs right inside it:",
    items=[
        ("Write niva in a Marimo cell", "a readable pipeline beside your Python and charts"),
        ("Reactive results", "change an input, the map and tables update — live"),
        ("Repeatable", "the notebook is the analysis; re-run it any time, same result"),
        ("Shareable", "hand someone the notebook, or export it — no setup ritual"),
    ])
_text(s, 1.25, 6.3, 11.2, 0.5,
      [[("No new tool to adopt — ", {"size": 18, "color": INK}),
        ("niva meets your people where they already work.", {"size": 18, "color": DEEP, "bold": True})]])
_footer(s, 7)

# ---- 8. ONE GRAMMAR, MANY SURFACES (graphic) ---------------------------
image_slide(prs, "One grammar, many surfaces", "surfaces.png", eyebrow="Flexible by design")
_footer(prs.slides[-1], 8)

# ---- 9. REPRODUCIBLE & PROVENANCE --------------------------------------
s = bullets_slide(
    prs, "Reproducible & shareable — by design", eyebrow="Trust the result",
    lead="The value decision-makers care about most: results you can re-run, audit, and defend.",
    items=[
        ("Pipelines are the record", "the one-line flow (or YAML) IS the documentation"),
        ("Provenance as a byproduct", "every step auto-recorded as lineage metadata"),
        ("Assess your inputs", "profile incoming data for quality before you trust it"),
        ("Version-controlled flows", "diff and review analysis like any other text"),
    ])
_footer(s, 9)

# ---- 10. EXPLORE ANYWHERE (graphic) ------------------------------------
image_slide(prs, "Produce once — explore it anywhere", "explore_anywhere.png",
            eyebrow="For every stakeholder")
_footer(prs.slides[-1], 10)

# ---- 11. WITHOUT NIVA (graphic) ----------------------------------------
image_slide(prs, "The workflow today", "without_niva.png", eyebrow="Without niva")
_footer(prs.slides[-1], 11)

# ---- 12. WITH NIVA (graphic) -------------------------------------------
image_slide(prs, "The workflow with niva", "with_niva.png", eyebrow="With niva")
_footer(prs.slides[-1], 12)

# ---- 13. THE REACH ------------------------------------------------------
s = content_slide(prs, "All of QGIS — in plain language", eyebrow="Thin, not limiting")
stats = [("769", "Processing algorithms"), ("406", "expression functions"), ("SQL", "SpatiaLite & PostGIS")]
cw, gap = 3.7, 0.3
x0 = (SW - (cw * 3 + gap * 2)) / 2
for i, (num, lab) in enumerate(stats):
    x = x0 + i * (cw + gap)
    card = _box(s, x, 2.1, cw, 2.3)
    _fill(card, LIGHT)
    _text(s, x, 2.35, cw, 1.2, num, size=58, color=DEEP, bold=True, align=PP_ALIGN.CENTER)
    _text(s, x, 3.7, cw, 0.6, lab, size=18, color=INK, align=PP_ALIGN.CENTER)
_text(s, 0.85, 4.95, 11.6, 1.4,
      [[("niva is a ", {"size": 20, "color": SAGE}),
        ("thin wrapper", {"size": 20, "color": INK, "bold": True}),
        (" — friendly verbs mapped onto QGIS's own algorithms, not a re-implementation of GIS. "
         "You get the full power of the world's leading open-source GIS, reachable in one readable line.",
         {"size": 20, "color": SAGE})]], line_spacing=1.2)
_footer(s, 13)

# ---- 14. WHY NOW --------------------------------------------------------
s = bullets_slide(
    prs, "Why niva, why now", eyebrow="The tailwinds",
    items=[
        ("QGIS is the de-facto open GIS", "ubiquitous, trusted, free — the foundation we build on"),
        ("Reproducible analysis is the norm", "Marimo and notebooks are how modern teams work"),
        ("The automation gap is real & unserved", "GUI analysts still can't automate without code"),
        ("Open and clean-room", "GPLv3, built on PyQGIS; no proprietary lock-in to fear"),
    ])
_footer(s, 14)

# ---- 15. VALUE / ROI ----------------------------------------------------
s = content_slide(prs, "What niva delivers", eyebrow="The value")
vals = [
    ("Speed", "answers the same day — no queue behind a programmer"),
    ("Capacity", "your existing analysts automate their own work"),
    ("Trust", "reproducible, provenance-tracked, auditable results"),
    ("Reach", "share to QGIS, Marimo, or static files — for anyone"),
    ("Lower cost", "less custom scripting, less specialized hiring"),
    ("No lock-in", "open source, on open foundations"),
]
cw, ch, gx, gy = 5.75, 1.35, 0.3, 0.25
x0 = (SW - (cw * 2 + gx)) / 2
for i, (h, b) in enumerate(vals):
    col, row = i % 2, i // 2
    x = x0 + col * (cw + gx)
    y = 1.95 + row * (ch + gy)
    card = _box(s, x, y, cw, ch)
    _fill(card, LIGHT)
    bar = _box(s, x, y, 0.14, ch)
    _fill(bar, TAN)
    _text(s, x + 0.35, y + 0.16, cw - 0.5, 0.5, h, size=20, color=DEEP, bold=True)
    _text(s, x + 0.35, y + 0.66, cw - 0.5, 0.6, b, size=15.5, color=SAGE)
_footer(s, 15)

# ---- 16. STATUS — WHAT'S BUILT -----------------------------------------
s = content_slide(prs, "What's working today", eyebrow="Status · built and shipping")
_text(s, 0.85, 1.55, 11.6, 0.5,
      [[("niva ships as a QGIS plugin (v0.27) and runs real analysis "
         "end-to-end, on QGIS's own algorithms:", {"size": 17, "color": SAGE})]])
status_cards = [
    ("Grammar & engine", DEEP, [
        "Readable lexer + parser → pipeline stages",
        "Pipe-chaining engine: layer handles, lineage",
        "PyQGIS backend — runs real geoprocessing",
    ]),
    ("Verbs & algorithms", TAN, [
        "~45 alias verbs + 12 built-ins (vector + raster)",
        "sql @conn — SELECT → layer, and server-side writes",
        "run — reach ANY of QGIS's 769 algorithms",
    ]),
    ("Data & cartography", DEEP, [
        "PostGIS & SpatiaLite — read · write · analyse",
        "project — repoint / build / templates · bookmarks",
        "style — apply / export .qml · .sld · .qlr",
    ]),
    ("Delivery & quality", TAN, [
        "QGIS plugin (Install from ZIP) · CLI · Python API",
        "Provenance journal · assess · catalog · notify / email",
        "270+ tests, live-QGIS + PostGIS CI · full docs",
    ]),
]
cw, ch, gx, gy = 5.8, 2.25, 0.25, 0.2
x0 = (SW - (cw * 2 + gx)) / 2
for i, (ctitle, accent, items) in enumerate(status_cards):
    col, row = i % 2, i // 2
    x = x0 + col * (cw + gx)
    y = 2.2 + row * (ch + gy)
    card = _box(s, x, y, cw, ch)
    _fill(card, LIGHT)
    bar = _box(s, x, y, cw, 0.16)
    _fill(bar, accent)
    _text(s, x + 0.3, y + 0.3, cw - 0.6, 0.4, ctitle, size=17, color=DEEP, bold=True)
    paras = [[("✓  ", {"color": DEEP, "bold": True, "size": 14}),
              (it, {"color": SAGE, "size": 14})] for it in items]
    _text(s, x + 0.3, y + 0.82, cw - 0.55, ch - 0.92, paras, line_spacing=1.3)
_footer(s, 16)

# ---- 17. ROADMAP — REMAINING MILESTONES --------------------------------
s = content_slide(prs, "The road ahead", eyebrow="Roadmap · shipped & next")
spine = _box(s, 1.62, 2.1, 0.04, 4.35)
_fill(spine, TAN)
milestones = [
    ("v0.27", "Shipped — available now",
     "~45 verbs + raster · run → 769 · PostGIS read·write·analyse · project / template / style · QGIS plugin · CLI / Python · full docs", "available now"),
    ("v1.0", "Stable release",
     "Grammar freeze (SemVer) · PyPI publish · worked Marimo–QGIS integration", None),
    ("v2.0", "Power features",
     "Named intermediates & variables · SQL-driven quality rules & constraints", None),
    ("v2.x", "Service mode",
     "Service / daemon mode · richer layout & symbology export", None),
]
ry, rh = 2.0, 1.02
for ver, mtitle, sub, tag in milestones:
    pill = _box(s, 1.0, ry, 1.25, 0.5)
    _fill(pill, DEEP)
    _text(s, 1.0, ry, 1.25, 0.5, ver, size=18, color=WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    head = [[(mtitle, {"size": 18, "color": INK, "bold": True})]]
    if tag:
        head[0].append(("   " + tag, {"size": 13, "color": PAIN, "bold": True, "italic": True}))
    _text(s, 2.6, ry - 0.04, 10.4, 0.4, head)
    _text(s, 2.6, ry + 0.4, 10.4, 0.5, sub, size=13.5, color=SAGE)
    ry += rh
_footer(s, 17)

# ---- 18. CTA ------------------------------------------------------------
s = prs.slides.add_slide(prs.slide_layouts[6])
_bg_white(s)
band = _box(s, 0, 0, SW, SH)
_fill(band, INK)
mark = os.path.join(LOGOS, "logo.png")
mw, mh = png_size(mark)
mh_in = 1.7
mw_in = mh_in * mw / mh
s.shapes.add_picture(mark, Inches((SW - mw_in) / 2), Inches(1.15), Inches(mw_in), Inches(mh_in))
_text(s, 1, 3.15, SW - 2, 0.9, "Let your analysts do the analysis.", size=34, color=WHITE,
      bold=True, align=PP_ALIGN.CENTER)
_text(s, 1, 4.05, SW - 2, 0.6,
      "Bring readable, reproducible geoprocessing to the people who need it.",
      size=19, color=TAN, align=PP_ALIGN.CENTER)
_text(s, 1, 5.0, SW - 2, 0.5, "Easy wins every time.", size=22, color=WHITE, italic=True,
      bold=True, align=PP_ALIGN.CENTER)
_text(s, 1, 6.4, SW - 2, 0.4, "github.com/johnzastrow/niva   ·   Let's run a pilot.",
      size=15, color=MIST, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("wrote", OUT, "—", len(prs.slides), "slides")
